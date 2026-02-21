#!/usr/bin/env python3
"""
Update a PowerPoint deck from PRESENTATION_PITCH_CONTENT.md.

Usage:
  python update_pitch_deck.py [--inspect] [--md PATH] [--pptx PATH] [--out PATH]

  --inspect     Print slide count and placeholder layout of the .pptx (no write).
  --md PATH     Path to PRESENTATION_PITCH_CONTENT.md (default: repo root).
  --pptx PATH   Path to source .pptx (default: ~/Downloads/Black Elegant and Modern Startup Pitch Deck Presentation.pptx).
  --out PATH    Output .pptx path (default: same dir as pptx, suffix _WAF_updated.pptx).

Requires: pip install python-pptx
"""

from __future__ import annotations

import argparse
import re
from pathlib import Path
from typing import NamedTuple


class SlideContent(NamedTuple):
    title: str
    body: str
    notes: str


def parse_pitch_markdown(md_path: Path) -> list[SlideContent]:
    """Parse PRESENTATION_PITCH_CONTENT.md into a list of (title, body, notes) per slide."""
    text = md_path.read_text(encoding="utf-8")

    # Split by ## Slide N: ...
    pattern = re.compile(r"^## Slide (\d+): (.+)$", re.MULTILINE)
    slides: list[SlideContent] = []
    last_end = 0
    matches = list(pattern.finditer(text))

    for i, m in enumerate(matches):
        slide_num = int(m.group(1))
        heading_title = m.group(2).strip()
        start = m.end()
        end = matches[i + 1].start() if i + 1 < len(matches) else len(text)
        block = text[start:end]

        # Extract body: lines that are bullets or numbered; skip mermaid blocks and speaker notes
        body_lines: list[str] = []
        notes_lines: list[str] = []
        in_mermaid = False
        for line in block.split("\n"):
            stripped = line.strip()
            if stripped.startswith("```"):
                in_mermaid = not in_mermaid
                continue
            if in_mermaid:
                continue
            if stripped.startswith("*Speaker note:*"):
                notes_lines.append(stripped.replace("*Speaker note:*", "").strip())
                continue
            if re.match(r"^[-*]\s+", stripped) or re.match(r"^\d+\.\s+", stripped):
                # Bullet or numbered: strip leading "- " or "1. " etc., then clean markdown
                clean = re.sub(r"^[-*]\s+", "", stripped)
                clean = re.sub(r"^\d+\.\s+", "", clean)
                clean = re.sub(r"\*\*(.+?)\*\*", r"\1", clean)  # bold -> plain
                clean = re.sub(r"\*(.+?)\*", r"\1", clean)      # italic -> plain
                clean = clean.replace("&lt;", "<").replace("&gt;", ">")
                body_lines.append(clean)
            elif stripped and not stripped.startswith("#") and slide_num == 1:
                # Slide 1 may have subtitle / your name on own line
                if "**Title:**" in stripped:
                    continue  # handled below
                if "**Subtitle:**" in stripped:
                    body_lines.append(stripped.replace("**Subtitle:**", "").strip())
                    continue
                if "**Your name**" in stripped or "Your name" in stripped:
                    body_lines.append(stripped)
                    continue

        body = "\n".join(body_lines).strip()
        notes = " ".join(notes_lines).strip()

        # Slide 1: title from "**Title:** ..." or heading; body = subtitle + your name
        if slide_num == 1:
            title_match = re.search(r"[-*]\s*\*\*Title:\*\*\s*(.+)", block)
            subtitle_match = re.search(r"[-*]\s*\*\*Subtitle:\*\*\s*(.+)", block)
            name_match = re.search(r"[-*]\s*\*\*Your name.*?\*\*\s*(.+)", block, re.IGNORECASE)
            title = title_match.group(1).strip() if title_match else heading_title
            body_parts = []
            if subtitle_match:
                body_parts.append(subtitle_match.group(1).strip())
            for line in block.split("\n"):
                s = line.strip()
                if "Your name" in s or "Course" in s or "Date" in s:
                    clean = re.sub(r"^[-*]\s*", "", s)
                    clean = re.sub(r"\*\*(.+?)\*\*", r"\1", clean)
                    if clean and clean not in body_parts:
                        body_parts.append(clean)
            body = "\n".join(body_parts) if body_parts else body
        else:
            title = heading_title

        slides.append(SlideContent(title=title, body=body, notes=notes))

    return slides


def inspect_pptx(pptx_path: Path) -> None:
    """Print slide count and placeholder layout for the given .pptx."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation(str(pptx_path))
    print(f"Slides: {len(prs.slides)}")
    for i, slide in enumerate(prs.slides):
        print(f"\n  Slide {i + 1}:")
        for shape in slide.shapes:
            if shape.is_placeholder:
                pf = shape.placeholder_format
                idx = getattr(pf, "idx", "?")
                ph_type = getattr(pf, "type", None)
                print(f"    idx={idx} type={ph_type} name={shape.name!r}")
            else:
                print(f"    (non-placeholder) name={shape.name!r}")


def _get_top(shape) -> float:
    """Vertical position for sorting (top of shape)."""
    try:
        return shape.top
    except Exception:
        return 0


def get_title_and_body_placeholders(slide) -> tuple[object | None, object | None]:
    """Return (title_shape, body_shape) for a slide. Prefer placeholders by idx 0 and 1; else first two text shapes by top."""
    title_ph = None
    body_ph = None
    for shape in slide.shapes:
        if shape.is_placeholder:
            try:
                idx = shape.placeholder_format.idx
                if idx == 0:
                    title_ph = shape
                elif idx == 1:
                    body_ph = shape
            except Exception:
                pass
            if shape.has_text_frame and title_ph is None and "title" in shape.name.lower():
                title_ph = shape
            if shape.has_text_frame and body_ph is None and ("content" in shape.name.lower() or "body" in shape.name.lower() or "text" in shape.name.lower()):
                body_ph = shape
    if title_ph and body_ph:
        return title_ph, body_ph
    # No placeholders (e.g. Canva-exported deck): use first two text-bearing shapes by top position
    text_shapes = [s for s in slide.shapes if s.has_text_frame and s.text.strip()]
    if not text_shapes:
        text_shapes = [s for s in slide.shapes if s.has_text_frame]
    text_shapes.sort(key=_get_top)
    if title_ph is None and text_shapes:
        title_ph = text_shapes[0]
    if body_ph is None and len(text_shapes) > 1:
        body_ph = text_shapes[1]
    return title_ph, body_ph


def set_placeholder_text(shape, text: str, max_chars: int = 32000) -> None:
    """Set text on a placeholder. Truncate if very long."""
    if not shape.has_text_frame:
        return
    text = (text or "")[:max_chars]
    try:
        shape.text = text
    except Exception:
        if shape.text_frame.paragraphs:
            p = shape.text_frame.paragraphs[0]
            if p.runs:
                p.runs[0].text = text
            else:
                p.text = text


def update_pptx(
    md_path: Path,
    pptx_path: Path,
    out_path: Path,
    slides_content: list[SlideContent],
    write_notes: bool = True,
) -> None:
    """Update the .pptx with slide content; add slides if needed."""
    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation(str(pptx_path))
    existing = len(prs.slides)
    # Layout for new slides: use layout of slide 1 (often title + content) or first content-like slide
    slide_layout = prs.slide_layouts[0]
    for idx, layout in enumerate(prs.slide_layouts):
        if "title" in layout.name.lower() and "content" in layout.name.lower():
            slide_layout = layout
            break
        if idx == 1:
            slide_layout = layout

    for i, content in enumerate(slides_content):
        if i < existing:
            slide = prs.slides[i]
        else:
            slide = prs.slides.add_slide(slide_layout)
        title_ph, body_ph = get_title_and_body_placeholders(slide)
        if title_ph and body_ph and title_ph is body_ph:
            # Single text box: put title and body together
            set_placeholder_text(title_ph, content.title + "\n\n" + content.body)
        else:
            if title_ph:
                set_placeholder_text(title_ph, content.title)
            if body_ph:
                set_placeholder_text(body_ph, content.body)
        if write_notes and content.notes and hasattr(slide, "notes_slide") and slide.notes_slide is not None:
            try:
                notes_slide = slide.notes_slide
                if notes_slide.notes_text_frame:
                    notes_slide.notes_text_frame.text = content.notes
            except Exception:
                pass

    prs.save(str(out_path))
    print(f"Saved {out_path} ({len(slides_content)} slides)")


def export_slide_mapping(slides_content: list[SlideContent], out_path: Path) -> None:
    """Write SLIDE_CONTENT_FOR_DECK.md for Canva / manual paste."""
    lines = [
        "# Slide content for deck (Canva / manual paste)",
        "",
        "Copy each section into the corresponding slide in Canva or PowerPoint.",
        "Diagrams for slides 5 and 6: see PRESENTATION_PITCH_CONTENT.md (Mermaid blocks) or redraw in editor.",
        "",
    ]
    for i, s in enumerate(slides_content, 1):
        lines.append("---")
        lines.append(f"## Slide {i}: {s.title}")
        lines.append("")
        lines.append("**Title (paste into title box):**")
        lines.append("")
        lines.append(s.title)
        lines.append("")
        lines.append("**Body (paste into content box):**")
        lines.append("")
        lines.append(s.body)
        if s.notes:
            lines.append("")
            lines.append("**Speaker note:**")
            lines.append(s.notes)
        lines.append("")
    out_path.write_text("\n".join(lines), encoding="utf-8")
    print(f"Wrote {out_path}")


def main() -> None:
    parser = argparse.ArgumentParser(description="Update pitch deck from PRESENTATION_PITCH_CONTENT.md")
    parser.add_argument("--inspect", action="store_true", help="Inspect .pptx placeholders only (no write)")
    parser.add_argument("--md", type=Path, default=None, help="Path to PRESENTATION_PITCH_CONTENT.md")
    parser.add_argument("--pptx", type=Path, default=None, help="Path to source .pptx")
    parser.add_argument("--out", type=Path, default=None, help="Output .pptx path")
    parser.add_argument("--export-md", type=Path, default=None, help="Export slide mapping to this .md file (e.g. SLIDE_CONTENT_FOR_DECK.md)")
    args = parser.parse_args()

    repo_root = Path(__file__).resolve().parent.parent
    md_path = args.md or repo_root / "PRESENTATION_PITCH_CONTENT.md"
    if not md_path.exists():
        raise SystemExit(f"Markdown not found: {md_path}")

    default_pptx = Path.home() / "Downloads" / "Black Elegant and Modern Startup Pitch Deck Presentation.pptx"
    pptx_path = args.pptx or default_pptx
    if not pptx_path.exists():
        raise SystemExit(f"PPTX not found: {pptx_path}")

    slides_content = parse_pitch_markdown(md_path)
    if len(slides_content) != 19:
        print(f"Warning: parsed {len(slides_content)} slides (expected 19)")

    if args.inspect:
        inspect_pptx(pptx_path)
        return

    if args.export_md:
        export_slide_mapping(slides_content, args.export_md)
        return

    if args.out:
        out_path = args.out
    else:
        out_path = pptx_path.parent / (pptx_path.stem + "_WAF_updated.pptx")

    update_pptx(md_path, pptx_path, out_path, slides_content)
    export_slide_mapping(slides_content, repo_root / "SLIDE_CONTENT_FOR_DECK.md")
    print("Done. For Canva, use SLIDE_CONTENT_FOR_DECK.md to copy content manually.")


if __name__ == "__main__":
    main()
